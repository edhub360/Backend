import fitz  # PyMuPDF
import docx
import pptx
import requests
from bs4 import BeautifulSoup
from youtube_transcript_api import YouTubeTranscriptApi
import io
import re
from urllib.parse import urlparse

async def extract_text_from_file(file):
    """Extract text from uploaded file (DEPRECATED - use extract_text_from_file_content)"""
    try:
        content = await file.read()
        return extract_text_from_file_content(content, file.filename)
    except Exception as e:
        error_msg = f"File processing error: {str(e)}"
        print(f"DEBUG: {error_msg}")
        return error_msg, {"filename": file.filename, "error": str(e)}

def extract_text_from_file_content(file_content: bytes, filename: str):
    """Extract text from file content bytes"""
    try:
        original_filename = filename
        filename_lower = filename.lower()
        metadata = {"filename": original_filename, "size": len(file_content)}
        
        print(f"DEBUG: Processing file: {filename}, size: {len(file_content)} bytes")
        
        if len(file_content) == 0:
            print("DEBUG: File is empty")
            return "File is empty", metadata
            
        # PDF Files
        if filename_lower.endswith(".pdf"):
            try:
                with fitz.open(stream=file_content, filetype="pdf") as doc:
                    text_parts = []
                    for page_num, page in enumerate(doc):
                        page_text = page.get_text()
                        if page_text.strip():
                            text_parts.append(page_text)
                    
                    text = "\n\n".join(text_parts)
                    print(f"DEBUG: PDF extracted {len(text)} characters from {len(doc)} pages")
                    
                    if len(text.strip()) == 0:
                        return "PDF appears to be empty or contains only images", metadata
                    
                    return text, metadata
                    
            except Exception as e:
                print(f"DEBUG: PDF extraction failed: {e}")
                return f"PDF processing failed: {str(e)}", metadata
        
        # Word Documents (.docx)
        elif filename_lower.endswith(".docx"):
            try:
                doc = docx.Document(io.BytesIO(file_content))
                text_parts = []
                
                for paragraph in doc.paragraphs:
                    if paragraph.text.strip():
                        text_parts.append(paragraph.text.strip())
                
                # Also extract text from tables
                for table in doc.tables:
                    for row in table.rows:
                        for cell in row.cells:
                            if cell.text.strip():
                                text_parts.append(cell.text.strip())
                
                text = "\n".join(text_parts)
                print(f"DEBUG: DOCX extracted {len(text)} characters")
                
                if len(text.strip()) == 0:
                    return "Word document appears to be empty", metadata
                
                return text, metadata
                
            except Exception as e:
                print(f"DEBUG: DOCX extraction failed: {e}")
                return f"DOCX processing failed: {str(e)}", metadata
        
        # PowerPoint Presentations (.pptx)
        elif filename_lower.endswith(".pptx"):
            try:
                ppt = pptx.Presentation(io.BytesIO(file_content))
                text_parts = []
                
                for slide_num, slide in enumerate(ppt.slides):
                    slide_texts = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_texts.append(shape.text.strip())
                    
                    if slide_texts:
                        text_parts.append(f"Slide {slide_num + 1}:\n" + "\n".join(slide_texts))
                
                text = "\n\n".join(text_parts)
                print(f"DEBUG: PPTX extracted {len(text)} characters from {len(ppt.slides)} slides")
                
                if len(text.strip()) == 0:
                    return "PowerPoint presentation appears to be empty", metadata
                
                return text, metadata
                
            except Exception as e:
                print(f"DEBUG: PPTX extraction failed: {e}")
                return f"PPTX processing failed: {str(e)}", metadata
        
        # Text Files (.txt)
        elif filename_lower.endswith(".txt"):
            try:
                # Try UTF-8 first, then fallback to other encodings
                encodings = ['utf-8', 'utf-16', 'iso-8859-1', 'cp1252']
                
                for encoding in encodings:
                    try:
                        text = file_content.decode(encoding)
                        print(f"DEBUG: TXT extracted {len(text)} characters using {encoding} encoding")
                        
                        if len(text.strip()) == 0:
                            return "Text file is empty", metadata
                        
                        return text, metadata
                    except UnicodeDecodeError:
                        continue
                
                # If all encodings fail
                return "Unable to decode text file - unsupported encoding", metadata
                
            except Exception as e:
                print(f"DEBUG: TXT extraction failed: {e}")
                return f"TXT processing failed: {str(e)}", metadata
        
        # Excel Files (.xlsx) - Basic support
        elif filename_lower.endswith(".xlsx"):
            try:
                import openpyxl
                wb = openpyxl.load_workbook(io.BytesIO(file_content))
                text_parts = []
                
                for sheet_name in wb.sheetnames:
                    sheet = wb[sheet_name]
                    sheet_texts = []
                    
                    for row in sheet.iter_rows(values_only=True):
                        row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                        if row_text.strip():
                            sheet_texts.append(row_text)
                    
                    if sheet_texts:
                        text_parts.append(f"Sheet: {sheet_name}\n" + "\n".join(sheet_texts))
                
                text = "\n\n".join(text_parts)
                print(f"DEBUG: XLSX extracted {len(text)} characters")
                
                if len(text.strip()) == 0:
                    return "Excel file appears to be empty", metadata
                
                return text, metadata
                
            except ImportError:
                return "Excel processing requires openpyxl package", metadata
            except Exception as e:
                print(f"DEBUG: XLSX extraction failed: {e}")
                return f"Excel processing failed: {str(e)}", metadata
        
        else:
            error_msg = f"Unsupported file type: {filename} (supported: .pdf, .docx, .pptx, .txt, .xlsx)"
            print(f"DEBUG: {error_msg}")
            return error_msg, metadata
            
    except Exception as e:
        error_msg = f"Unexpected file processing error: {str(e)}"
        print(f"DEBUG: {error_msg}")
        return error_msg, {"filename": filename, "error": str(e)}

async def extract_from_url(url):
    """Extract text from a website URL"""
    try:
        print(f"DEBUG: Extracting from URL: {url}")
        
        # Validate URL
        parsed_url = urlparse(url)
        if not parsed_url.scheme or not parsed_url.netloc:
            return "Invalid URL format", {"url": url}
        
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
        }
        
        response = requests.get(url, headers=headers, timeout=10)
        response.raise_for_status()
        
        soup = BeautifulSoup(response.content, 'html.parser')
        
        # Remove script and style elements
        for script in soup(["script", "style", "nav", "header", "footer"]):
            script.decompose()
        
        # Extract text from common content areas
        content_selectors = [
            'article', 'main', '.content', '#content', '.post', '.article'
        ]
        
        text_parts = []
        
        # Try to find main content
        for selector in content_selectors:
            elements = soup.select(selector)
            if elements:
                for element in elements:
                    text = element.get_text(separator='\n', strip=True)
                    if len(text) > 100:  # Only include substantial content
                        text_parts.append(text)
                break
        
        # If no specific content found, get all text
        if not text_parts:
            text = soup.get_text(separator='\n', strip=True)
            # Clean up excessive whitespace
            text = re.sub(r'\n\s*\n', '\n\n', text)
            text_parts.append(text)
        
        extracted_text = '\n\n'.join(text_parts)
        
        # Get page title
        title = soup.title.string if soup.title else "No title"
        
        metadata = {
            "url": url,
            "title": title.strip(),
            "length": len(extracted_text)
        }
        
        print(f"DEBUG: URL extracted {len(extracted_text)} characters")
        
        if len(extracted_text.strip()) == 0:
            return "No readable content found on the webpage", metadata
        
        return extracted_text, metadata
        
    except requests.exceptions.RequestException as e:
        error_msg = f"Failed to fetch URL: {str(e)}"
        print(f"DEBUG: {error_msg}")
        return error_msg, {"url": url, "error": str(e)}
    except Exception as e:
        error_msg = f"URL processing error: {str(e)}"
        print(f"DEBUG: {error_msg}")
        return error_msg, {"url": url, "error": str(e)}

async def extract_from_youtube(youtube_url):
    """Extract transcript from YouTube video"""
    try:
        print(f"DEBUG: Extracting from YouTube: {youtube_url}")
        
        # Extract video ID from URL
        video_id = None
        patterns = [
            r'(?:youtube\.com\/watch\?v=|youtu\.be\/|youtube\.com\/embed\/)([a-zA-Z0-9_-]{11})',
        ]
        
        for pattern in patterns:
            match = re.search(pattern, youtube_url)
            if match:
                video_id = match.group(1)
                break
        
        if not video_id:
            return "Invalid YouTube URL format", {"url": youtube_url}
        
        # Get transcript
        transcript_list = YouTubeTranscriptApi.get_transcript(video_id)
        
        # Combine transcript entries
        transcript_text = ' '.join([entry['text'] for entry in transcript_list])
        
        # Clean up the text
        transcript_text = re.sub(r'\s+', ' ', transcript_text).strip()
        
        metadata = {
            "url": youtube_url,
            "video_id": video_id,
            "transcript_entries": len(transcript_list),
            "length": len(transcript_text)
        }
        
        print(f"DEBUG: YouTube extracted {len(transcript_text)} characters from {len(transcript_list)} transcript entries")
        
        if len(transcript_text.strip()) == 0:
            return "No transcript available for this video", metadata
        
        return transcript_text, metadata
        
    except Exception as e:
        error_msg = f"YouTube processing error: {str(e)}"
        print(f"DEBUG: {error_msg}")
        return error_msg, {"url": youtube_url, "error": str(e)}
